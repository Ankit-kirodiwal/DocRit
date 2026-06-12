import os
import sys
import argparse
import subprocess
import glob
import fitz # PyMuPDF
import pandas as pd
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pdf2docx import Converter
import re

# ==========================================
# Monkey-patches to optimize pdf2docx layout
# ==========================================

# 1. Spans.restore patch to keep whitespace-only spans (prevent word-merging)
from pdf2docx.text.Spans import Spans
from pdf2docx.text.TextSpan import TextSpan
from pdf2docx.image.ImageSpan import ImageSpan

def _patched_spans_restore(self, raws: list):
    for raw_span in raws:
        if 'image' in raw_span:
            span = ImageSpan(raw_span)
        else:
            span = TextSpan(raw_span)
            if not span.text and not span.style:
                span = None
        self.append(span)
    return self

Spans.restore = _patched_spans_restore

# 2. RawPage.parse_section patch to refine multi-column detection thresholds
from pdf2docx.page.RawPage import RawPage
from pdf2docx.common.Collection import Collection

def _patched_parse_section(self, **settings):
    X0, Y0, X1, _ = self.working_bbox

    # collect all blocks (line level) and shapes
    elements = Collection()
    elements.extend(self.blocks)
    elements.extend(self.shapes.text_style_shapes)
    if not elements: return []

    # to create section with collected lines
    lines = Collection()
    sections = []
    def close_section(num_col, elements, y_ref):
        # append to last section if both single column
        if sections and sections[-1].num_cols==num_col==1:
            column = sections[-1][0]
            column.union_bbox(elements)
            column.add_elements(elements)
        # otherwise, create new section
        else:
            section = self._create_section(num_col, elements, (X0, X1), y_ref)
            if section:
                sections.append(section)

    # check section row by row
    pre_num_col = 1
    y_ref = Y0 # to calculate v-distance between sections
    for row in elements.group_by_rows():
        # check column col by col
        cols = row.group_by_columns()
        current_num_col = len(cols)

        # column check:
        # consider 2-cols only
        if current_num_col>2:
            current_num_col = 1

        # the width of two columns shouldn't have significant difference
        elif current_num_col==2:
            u0, v0, u1, v1 = cols[0].bbox
            m0, n0, m1, n1 = cols[1].bbox
            x0 = (u1+m0)/2.0
            c1, c2 = x0-X0, X1-x0 # column width
            w1, w2 = u1-u0, m1-m0 # line width
            f = 2.5 # patched from 2.0 to support short text in columns (e.g. Abstract)
            if not 1/f<=c1/c2<=f or w1/c1<0.1 or w2/c2<0.1: # patched from 0.33 to 0.1
                current_num_col = 1

        # process exceptions
        if pre_num_col==2 and current_num_col==1:
            # though current row has one single column, it might have another virtual
            # and empty column. If so, it should be counted as 2-cols
            cols = lines.group_by_columns()
            pos = cols[0].bbox[2]
            if row.bbox[2]<=pos or row.bbox[0]>pos:
                current_num_col = 2

            # pre_num_col!=current_num_col => to close section with collected lines,
            # before that, further check the height of collected lines
            else:
                x0, y0, x1, y1 = lines.bbox
                if y1-y0<settings['min_section_height']:
                    pre_num_col = 1

        elif pre_num_col==2 and current_num_col==2:
            # though both 2-cols, they don't align with each other
            combine = Collection(lines)
            combine.extend(row)
            if len(combine.group_by_columns(sorted=False))==1: current_num_col = 1

        # finalize pre-section if different from the column count of previous section
        if current_num_col!=pre_num_col:
            # process pre-section
            close_section(pre_num_col, lines, y_ref)
            if sections:
                y_ref = sections[-1][-1].bbox[3]

            # start potential new section
            lines = Collection(row)
            pre_num_col = current_num_col

        # otherwise, collect current lines for further processing
        else:
            lines.extend(row)

    # don't forget the final section
    close_section(current_num_col, lines, y_ref)

    return sections

RawPage.parse_section = _patched_parse_section


# Find Tesseract executable path on Windows
def find_tesseract():
    # Try PATH first
    try:
        subprocess.run(["tesseract", "--version"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        return "tesseract"
    except FileNotFoundError:
        pass
    
    # Check common installation locations
    common_paths = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        os.path.join(os.environ.get("LOCALAPPDATA", ""), r"Programs\Tesseract-OCR\tesseract.exe"),
        r"D:\Program Files\Tesseract-OCR\tesseract.exe"
    ]
    for p in common_paths:
        if os.path.exists(p):
            return p
    return None

def convert_pdf_to_docx(pdf_path, docx_path):
    print(f"Converting PDF to DOCX: {pdf_path} -> {docx_path}")
    
    # Preprocess the PDF first to replace complex vector graphs/charts with clean raster images
    temp_pdf_path = f"temp_preprocessed_{os.getpid()}.pdf"
    pdf_to_convert = pdf_path
    
    try:
        doc = fitz.open(pdf_path)
        has_changes = False
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # Identify drawing blocks (clustered vector drawings)
            drawings = page.get_drawings()
            drawing_rects = []
            for d in drawings:
                r = d.get("rect")
                if r and not r.is_empty:
                    # Ignore large structural lines/frames
                    if r.width > 350 or r.height > 350:
                        continue
                    drawing_rects.append(r)
                    
            # Cluster drawing rects
            clusters = []
            for r in drawing_rects:
                merged = False
                for idx, c in enumerate(clusters):
                    test_c = fitz.Rect(c.x0 - 15, c.y0 - 15, c.x1 + 15, c.y1 + 15)
                    if test_c.intersects(r):
                        clusters[idx] = c | r
                        merged = True
                        break
                if not merged:
                    clusters.append(r)
                    
            changed = True
            while changed:
                changed = False
                i = 0
                while i < len(clusters):
                    j = i + 1
                    while j < len(clusters):
                        test_c = fitz.Rect(clusters[i].x0 - 15, clusters[i].y0 - 15, clusters[i].x1 + 15, clusters[i].y1 + 15)
                        if test_c.intersects(clusters[j]):
                            clusters[i] = clusters[i] | clusters[j]
                            clusters.pop(j)
                            changed = True
                        else:
                            j += 1
                    i += 1
                    
            # Get page dimensions and words to protect layout structures
            words = page.get_text("words")
            page_height = page.rect.height

            # Detect tables on the page to avoid rasterizing them
            table_rects = []
            try:
                tables = page.find_tables()
                if tables:
                    for t in tables:
                        if hasattr(t, "bbox"):
                            table_rects.append(fitz.Rect(t.bbox))
            except Exception as table_err:
                print(f"Table detection skipped: {table_err}")

            # Filter clusters to get valid graphic blocks
            graph_rects = []
            for c in clusters:
                c = c & page.rect
                if c.width > 20 and c.height > 20:
                    # 1. Skip header and footer zones (protects headers, footers, and page numbers)
                    if c.y1 < 75 or c.y0 > (page_height - 75):
                        continue

                    # 2. Avoid rasterizing any graphics that overlap with tables
                    is_table = False
                    for t_rect in table_rects:
                        if c.intersects(t_rect):
                            is_table = True
                            break
                    if is_table:
                        continue

                    # 3. Avoid rasterizing background shapes/watermarks that overlap with body text
                    # (Graphs only have short labels; body text has many characters)
                    overlapping_words = [w for w in words if fitz.Rect(w[:4]).intersects(c)]
                    total_char_len = sum(len(w[4]) for w in overlapping_words)
                    if total_char_len > 120:
                        continue

                    paths_in_cluster = [r for r in drawing_rects if c.contains(r)]
                    if len(paths_in_cluster) >= 5:
                        graph_rects.append(c)
                        
            # Replace graphs with high-quality PNGs
            for idx, g_rect in enumerate(graph_rects):
                pix = page.get_pixmap(clip=g_rect, dpi=300)
                graph_img_path = f"temp_pre_graph_{page_num}_{idx}_{os.getpid()}.png"
                pix.save(graph_img_path)
                
                # Redact the vector drawings in this area
                page.add_redact_annot(g_rect, fill=(1, 1, 1)) # fill with white
                page.apply_redactions()
                
                # Insert the rasterized image back in
                page.insert_image(g_rect, filename=graph_img_path)
                has_changes = True
                
                try:
                    os.remove(graph_img_path)
                except Exception:
                    pass
                    
        if has_changes:
            doc.save(temp_pdf_path)
            pdf_to_convert = temp_pdf_path
        doc.close()
        
    except Exception as err:
        print(f"Failed to preprocess PDF graphs: {err}. Proceeding with original PDF.")
        pdf_to_convert = pdf_path
        
    try:
        cv = Converter(pdf_to_convert)
        # Use default settings
        cv.convert(docx_path, start=0, end=None)
        cv.close()
    finally:
        if pdf_to_convert == temp_pdf_path and os.path.exists(temp_pdf_path):
            try:
                os.remove(temp_pdf_path)
            except Exception:
                pass

def convert_pdf_to_xlsx(pdf_path, xlsx_path):
    print(f"Converting PDF to XLSX: {pdf_path} -> {xlsx_path}")
    import pdfplumber
    from openpyxl.utils import get_column_letter
    from openpyxl.styles import Alignment
    import re

    # Detect if PDF has table borders (grid lines)
    is_bordered = False
    try:
        with pdfplumber.open(pdf_path) as pdf:
            if pdf.pages:
                # If first page has horizontal/vertical lines that form a grid, we treat it as bordered
                first_page = pdf.pages[0]
                lines_count = len(first_page.lines) + len(first_page.rects)
                if lines_count > 10:
                    is_bordered = True
    except Exception as e:
        print(f"Failed to detect table border type: {e}")

    # Try extracting tables with tabula only if PDF has borders
    dfs = []
    if is_bordered:
        try:
            import tabula
            print("Running Tabula table extraction for bordered PDF...")
            dfs = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
        except Exception as e:
            print(f"Tabula failed or not installed: {e}. Falling back to pdfplumber...")
            dfs = []

    with pd.ExcelWriter(xlsx_path, engine='openpyxl') as writer:
        if dfs and len(dfs) > 0:
            try:
                # Concatenate all tables into a single DataFrame to put all data on a single sheet
                combined_df = pd.concat(dfs, ignore_index=True)
                combined_df.to_excel(writer, sheet_name="Extracted Tables", index=False)
                print(f"Successfully wrote combined tables extracted via Tabula.")
            except Exception as concat_err:
                print(f"Failed to concatenate Tabula tables: {concat_err}. Writing individually...")
                for i, df in enumerate(dfs):
                    sheet_name = f"Table {i+1}"[:30]
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
        else:
            # Fallback to pdfplumber table extraction
            print("Running pdfplumber table extraction...")
            all_rows = []
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            all_rows.extend(table)
            
            if all_rows:
                df = pd.DataFrame(all_rows)
                df.to_excel(writer, sheet_name="Extracted Tables", index=False, header=False)
                tables_found = True
            else:
                tables_found = False
            
            if not tables_found:
                # Text extraction fallback
                print("No tables found. Performing text layout extraction...")
                
                def split_merged_chunk(text, x0, x1):
                    # Check if age + date is merged: e.g. "3215/10/2017" or "32 15/10/2017"
                    date_match = re.match(r"^(\d+)\s*(\d{2}/\d{2}/\d{4})$", text)
                    if date_match:
                        age_str = date_match.group(1)
                        date_str = date_match.group(2)
                        total_len = len(text)
                        char_width = (x1 - x0) / total_len
                        chunk1 = {'text': age_str, 'x0': x0, 'x1': x0 + char_width * len(age_str)}
                        chunk2 = {'text': date_str, 'x0': chunk1['x1'], 'x1': x1}
                        return [chunk1, chunk2]
                        
                    # Check if serial number + name is merged: e.g. "1Dulce" or "1 Dulce"
                    name_match = re.match(r"^(\d+)\s+([A-Z][a-zA-Z]*)$", text)
                    if not name_match:
                        name_match = re.match(r"^(\d+)([A-Z][a-zA-Z]*)$", text)
                    if name_match:
                        serial_str = name_match.group(1)
                        name_str = name_match.group(2)
                        total_len = len(text)
                        char_width = (x1 - x0) / total_len
                        # Account for space offset if space is present
                        space_offset = 0
                        if ' ' in text:
                            space_offset = char_width
                        chunk1 = {'text': serial_str, 'x0': x0, 'x1': x0 + char_width * len(serial_str)}
                        chunk2 = {'text': name_str, 'x0': chunk1['x1'] + space_offset, 'x1': x1}
                        return [chunk1, chunk2]
                        
                    return [{'text': text, 'x0': x0, 'x1': x1}]

                lines = []
                gapped_starts = []
                header_row = None
                with pdfplumber.open(pdf_path) as pdf:
                    for page in pdf.pages:
                        words = page.extract_words()
                        if not words:
                            continue
                            
                        # Group words by top coordinate with tolerance of 3 points
                        lines_grouped = {}
                        for w in words:
                            top = w['top']
                            found = False
                            for t in lines_grouped:
                                if abs(t - top) < 3.0:
                                    lines_grouped[t].append(w)
                                    found = True
                                    break
                            if not found:
                                lines_grouped[top] = [w]
                                
                        sorted_tops = sorted(lines_grouped.keys())
                        
                        # 1. Pre-merge words on each line into continuous text chunks (golden gap < 3.0 points)
                        page_lines = []
                        for top in sorted_tops:
                            line_words = sorted(lines_grouped[top], key=lambda x: x['x0'])
                            chunks = []
                            if line_words:
                                curr_text = line_words[0]['text']
                                start_x = line_words[0]['x0']
                                prev_x1 = line_words[0]['x1']
                                
                                for w in line_words[1:]:
                                    gap = w['x0'] - prev_x1
                                    if gap < 3.0:
                                        curr_text += " " + w['text']
                                    else:
                                        splits = split_merged_chunk(curr_text, start_x, prev_x1)
                                        chunks.extend(splits)
                                        curr_text = w['text']
                                        start_x = w['x0']
                                    prev_x1 = w['x1']
                                splits = split_merged_chunk(curr_text, start_x, prev_x1)
                                chunks.extend(splits)
                            
                            # Final pass to ensure all merged tokens are split
                            final_chunks = []
                            for ch in chunks:
                                splits = split_merged_chunk(ch['text'], ch['x0'], ch['x1'])
                                final_chunks.extend(splits)
                                
                            page_lines.append((top, final_chunks))
                            for ch in final_chunks:
                                gapped_starts.append(ch['x0'])
                        
                        lines.append(page_lines)
                                
                # 2. Cluster chunk starts globally to find column coordinates
                clusters = []
                for x in sorted(gapped_starts):
                    added = False
                    for c in clusters:
                        if abs(c['values'][0] - x) < 2.0:
                            c['values'].append(x)
                            c['center'] = sum(c['values']) / len(c['values'])
                            added = True
                            break
                    if not added:
                        clusters.append({'center': x, 'values': [x]})
                        
                valid_cols = sorted([c['center'] for c in clusters])
                
                if not valid_cols:
                    valid_cols = [0.0]
                    
                # Compute boundaries
                boundaries = []
                for idx in range(len(valid_cols) - 1):
                    boundaries.append((valid_cols[idx] + valid_cols[idx+1]) / 2.0)
                    
                # 3. Assign chunks to cells using boundaries and perform header deduplication
                flat_rows = []
                for page_lines in lines:
                    for top, chunks in page_lines:
                        # Skip common page/sheet titles
                        if len(chunks) == 1 and chunks[0]['text'].strip().lower() in ['sheet1', 'sheet 1']:
                            continue
                            
                        cells = [""] * len(valid_cols)
                        for ch in chunks:
                            col_idx = len(valid_cols) - 1
                            for idx, b in enumerate(boundaries):
                                if ch['x0'] < b:
                                    col_idx = idx
                                    break
                            # Append instead of overwriting to prevent data loss
                            cells[col_idx] = f"{cells[col_idx]} {ch['text']}".strip() if cells[col_idx] else ch['text']
                            
                        if not header_row and any(cells):
                            header_row = cells
                            flat_rows.append(cells)
                        elif header_row and cells == header_row:
                            continue
                        else:
                            flat_rows.append(cells)
                            
                # 4. Merge adjacent mutually exclusive columns
                num_cols = len(valid_cols)
                j = 0
                while j < num_cols - 1:
                    can_merge = True
                    for row in flat_rows:
                        if row[j] != "" and row[j+1] != "":
                            can_merge = False
                            break
                    
                    if can_merge and (valid_cols[j+1] - valid_cols[j]) < 150.0:
                        for row in flat_rows:
                            if row[j+1] != "":
                                row[j] = row[j+1]
                        for row in flat_rows:
                            row.pop(j+1)
                        valid_cols.pop(j+1)
                        num_cols -= 1
                    else:
                        j += 1
                        
                if not flat_rows:
                    flat_rows = [["No text or tables found in PDF"]]
                df = pd.DataFrame(flat_rows)
                df.to_excel(writer, sheet_name="Extracted Text", index=False, header=False)

        # Auto-adjust column widths and apply numeric right indentation with spacing
        for sheet_name in writer.sheets:
            ws = writer.sheets[sheet_name]
            for col in ws.columns:
                max_len = 0
                for cell in col:
                    if cell.value is not None:
                        max_len = max(max_len, len(str(cell.value)))
                        # Apply right-alignment with indent=1 to numbers
                        is_numeric = isinstance(cell.value, (int, float))
                        if not is_numeric and isinstance(cell.value, str):
                            cleaned = cell.value.strip().replace(',', '')
                            if re.match(r'^-?\d+(\.\d+)?$', cleaned):
                                is_numeric = True
                        
                        if is_numeric:
                            cell.alignment = Alignment(wrap_text=True, vertical='center', horizontal='right', indent=1)
                        else:
                            cell.alignment = Alignment(wrap_text=True, vertical='center', horizontal='left')
                col_letter = get_column_letter(col[0].column)
                ws.column_dimensions[col_letter].width = min(max(max_len + 5, 12), 30)

FONT_MAPPING = {
    "arialmt": "Arial",
    "arial": "Arial",
    "calibri": "Calibri",
    "timesnewromanpsmt": "Times New Roman",
    "timesnewroman": "Times New Roman",
    "times new roman": "Times New Roman",
    "courier": "Courier New",
    "couriernew": "Courier New",
    "helvetica": "Arial",
    "georgia": "Georgia",
    "verdana": "Verdana",
    "cambria": "Cambria",
    "tahoma": "Tahoma",
    "trebuchet": "Trebuchet MS",
    "lucida": "Lucida Sans"
}

def map_font(font_name):
    if not font_name:
        return "Arial"
    normalized = font_name.lower().replace("-", "").replace("bold", "").replace("italic", "").strip()
    return FONT_MAPPING.get(normalized, "Arial")

def convert_pdf_to_pptx(pdf_path, pptx_path):
    print(f"Converting PDF to PPTX (ultra-high-fidelity native): {pdf_path} -> {pptx_path}")
    import math
    
    prs = Presentation()
    doc = fitz.open(pdf_path)
    
    if len(doc) > 0:
        first_page = doc.load_page(0)
        rect = first_page.rect
        prs.slide_width = Inches(rect.width / 72.0)
        prs.slide_height = Inches(rect.height / 72.0)

    # Iterate through pages, adding background image and transparent editable text boxes
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        page_dict = page.get_text("dict")
        
        # Identify image blocks (type 1)
        image_blocks = []
        for block in page_dict.get("blocks", []):
            if block.get("type") == 1:
                image_blocks.append(block)
                
        # Identify drawing blocks (clustered vector drawings)
        drawings = page.get_drawings()
        drawing_rects = []
        for d in drawings:
            r = d.get("rect")
            if r and not r.is_empty:
                if r.width > page.rect.width * 0.95 and r.height > page.rect.height * 0.95:
                    continue
                drawing_rects.append(r)
                
        # Cluster drawing rects
        clusters = []
        for r in drawing_rects:
            merged = False
            for idx, c in enumerate(clusters):
                test_c = fitz.Rect(c.x0 - 15, c.y0 - 15, c.x1 + 15, c.y1 + 15)
                if test_c.intersects(r):
                    clusters[idx] = c | r
                    merged = True
                    break
            if not merged:
                clusters.append(r)
                
        changed = True
        while changed:
            changed = False
            i = 0
            while i < len(clusters):
                j = i + 1
                while j < len(clusters):
                    test_c = fitz.Rect(clusters[i].x0 - 15, clusters[i].y0 - 15, clusters[i].x1 + 15, clusters[i].y1 + 15)
                    if test_c.intersects(clusters[j]):
                        clusters[i] = clusters[i] | clusters[j]
                        clusters.pop(j)
                        changed = True
                    else:
                        j += 1
                i += 1
                
        # Filter clusters to get valid graphic blocks
        graph_rects = []
        for c in clusters:
            c = c & page.rect
            if c.width > 20 and c.height > 20:
                paths_in_cluster = [r for r in drawing_rects if c.contains(r)]
                if len(paths_in_cluster) >= 5:
                    graph_rects.append(c)

        # Create a temporary copy to redact text, images, and graphs, then render background
        temp_doc = fitz.open(pdf_path)
        temp_page = temp_doc.load_page(page_num)
        
        # Redact text
        words = temp_page.get_text("words")
        for w in words:
            rect = fitz.Rect(w[:4])
            temp_page.add_redact_annot(rect, fill=None)
            
        # Redact images
        for img in image_blocks:
            rect = fitz.Rect(img["bbox"])
            temp_page.add_redact_annot(rect, fill=None)
            
        # Redact graphs
        for g_rect in graph_rects:
            temp_page.add_redact_annot(g_rect, fill=None)
            
        temp_page.apply_redactions()
        
        # Render clean background image
        pix = temp_page.get_pixmap(dpi=150)
        bg_img_path = f"temp_bg_{page_num}_{os.getpid()}.png"
        pix.save(bg_img_path)
        
        # Add a blank slide
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Insert background image spanning the entire slide
        slide.shapes.add_picture(bg_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        try:
            os.remove(bg_img_path)
        except Exception:
            pass
        temp_doc.close()
        
        # Insert standard images as separate movable shapes
        for idx, img in enumerate(image_blocks):
            bbox = img["bbox"]
            left = Inches(bbox[0] / 72.0)
            top = Inches(bbox[1] / 72.0)
            width = Inches((bbox[2] - bbox[0]) / 72.0)
            height = Inches((bbox[3] - bbox[1]) / 72.0)
            
            img_bytes = img["image"]
            img_ext = img["ext"]
            img_path = f"temp_img_{page_num}_{idx}_{os.getpid()}.{img_ext}"
            with open(img_path, "wb") as f:
                f.write(img_bytes)
                
            try:
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)
            except Exception as e:
                print(f"Failed to insert image shape: {e}")
                
            try:
                os.remove(img_path)
            except Exception:
                pass
                
        # Render and insert graph drawing blocks as separate movable shapes
        for idx, g_rect in enumerate(graph_rects):
            pix = page.get_pixmap(clip=g_rect, dpi=300)
            graph_path = f"temp_graph_{page_num}_{idx}_{os.getpid()}.png"
            pix.save(graph_path)
            
            left = Inches(g_rect.x0 / 72.0)
            top = Inches(g_rect.y0 / 72.0)
            width = Inches(g_rect.width / 72.0)
            height = Inches(g_rect.height / 72.0)
            
            try:
                slide.shapes.add_picture(graph_path, left, top, width=width, height=height)
            except Exception as e:
                print(f"Failed to insert graph shape: {e}")
                
            try:
                os.remove(graph_path)
            except Exception:
                pass

        # Place text boxes on top of the background image
        for block in page_dict.get("blocks", []):
            if block.get("type") == 0:  # Text block
                for line in block.get("lines", []):
                    l_bbox = line.get("bbox")
                    spans = [s for s in line.get("spans", []) if s.get("text", "").strip()]
                    if not spans:
                        continue

                    max_font_size = max([s.get("size", 12) for s in spans])

                    # Calculate inches coordinates
                    left = Inches(l_bbox[0] / 72.0)
                    top = Inches(l_bbox[1] / 72.0)
                    width = Inches(max(l_bbox[2] - l_bbox[0], max_font_size * 2.0) / 72.0)
                    height = Inches(max(l_bbox[3] - l_bbox[1], max_font_size * 1.5) / 72.0)

                    # Add transparent text box
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = False
                    tf.margin_left = Inches(0.0)
                    tf.margin_right = Inches(0.0)
                    tf.margin_top = Inches(0.0)
                    tf.margin_bottom = Inches(0.0)

                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.LEFT

                    # Map rotation if the text line is tilted
                    dx, dy = line.get("dir", (1.0, 0.0))
                    angle = math.degrees(math.atan2(dy, dx))
                    if angle < 0:
                        angle += 360
                    if abs(angle) > 0.01:
                        txBox.rotation = angle

                    # Add text spans with matching styles
                    for span in spans:
                        run = p.add_run()
                        run.text = span.get("text", "")
                        run.font.name = map_font(span.get("font"))
                        run.font.size = Pt(span.get("size", 12))

                        # Color
                        color_val = span.get("color", 0)
                        r = (color_val >> 16) & 255
                        g = (color_val >> 8) & 255
                        b = color_val & 255
                        run.font.color.rgb = RGBColor(r, g, b)

                        # Flags (Bold / Italic)
                        flags = span.get("flags", 0)
                        if flags & 2:
                            run.font.italic = True
                        if flags & 16:
                            run.font.bold = True

    prs.save(pptx_path)
    doc.close()
    print("PPTX conversion complete.")


def perform_ocr(input_path, output_path, output_type):
    tesseract_cmd = find_tesseract()
    if not tesseract_cmd:
        print("CRITICAL ERROR: Tesseract OCR executable not found on the system.", file=sys.stderr)
        print("Please install Tesseract OCR and verify it is added to your PATH environment variable.", file=sys.stderr)
        sys.exit(2)
        
    print(f"Running Advanced OCR on {input_path} using Tesseract: {tesseract_cmd} for output format: {output_type}")
    
    # Mathematical character post-processing cleanup function
    def clean_ocr(text):
        # 1. Correct projection operator π (pi) misrecognized as n or m
        text = re.sub(r'\b([nm])([A-Z][a-z]+)\b', lambda m: 'π' + m.group(2), text)
        text = re.sub(r'\b([nm])([A-Z][a-z]+)\(', lambda m: 'π' + m.group(2) + '(', text)
        # 2. Correct >= (greater than or equal) misrecognized as Sigma or duplicated digit
        text = re.sub(r'(\b[a-zA-Z_][a-zA-Z0-9_]*\s+)[ΣΣ]\s*(\d+)', r'\1>= \2', text)
        text = re.sub(r'(\b[a-zA-Z_][a-zA-Z0-9_]*\s+)(\d+)\s+\2\b', r'\1>= \2', text)
        # 3. Clean up other common relational algebra math symbol misrecognitions
        text = re.sub(r'\b(Pd|bd|pa|Ρᾷ)\b', '⋈', text)
        return text

    tessdata_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "tessdata")
    system_tessdata = os.path.join(os.path.dirname(tesseract_cmd), "tessdata")
    if os.path.exists(system_tessdata):
        system_configs = os.path.join(system_tessdata, "configs")
        local_configs = os.path.join(tessdata_dir, "configs")
        if os.path.exists(system_configs) and not os.path.exists(local_configs):
            import shutil
            try:
                shutil.copytree(system_configs, local_configs, dirs_exist_ok=True)
            except Exception as e:
                print(f"Warning: Failed to copy configs folder: {e}")

    langs = "eng"
    extra_args = []
    if os.path.exists(os.path.join(tessdata_dir, "grc.traineddata")):
        langs = "eng+grc"
        extra_args = ["--tessdata-dir", tessdata_dir, "-l", langs]

    is_pdf = input_path.lower().endswith('.pdf')
    temp_files = []
    
    # Extract page images
    if is_pdf:
        doc = fitz.open(input_path)
        for i in range(len(doc)):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=200)
            img_path = f"temp_ocr_p{i}_{os.getpid()}.png"
            pix.save(img_path)
            temp_files.append((i, img_path))
        doc.close()
    else:
        temp_files.append((0, input_path))

    # Perform OCR based on output type
    if output_type == 'pdf':
        pdf_merger = fitz.open()
        for idx, img_path in temp_files:
            pdf_base = f"temp_ocr_out_p{idx}_{os.getpid()}"
            cmd = [tesseract_cmd] + extra_args + [img_path, pdf_base, "pdf"]
            subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
            page_pdf = pdf_base + ".pdf"
            if os.path.exists(page_pdf):
                page_doc = fitz.open(page_pdf)
                pdf_merger.insert_pdf(page_doc)
                page_doc.close()
                os.remove(page_pdf)
        pdf_merger.save(output_path)
        pdf_merger.close()
        
    elif output_type == 'docx':
        # Create searchable PDF, then convert to docx
        temp_searchable = f"temp_searchable_{os.getpid()}.pdf"
        pdf_merger = fitz.open()
        for idx, img_path in temp_files:
            pdf_base = f"temp_ocr_out_p{idx}_{os.getpid()}"
            cmd = [tesseract_cmd] + extra_args + [img_path, pdf_base, "pdf"]
            subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
            page_pdf = pdf_base + ".pdf"
            if os.path.exists(page_pdf):
                page_doc = fitz.open(page_pdf)
                pdf_merger.insert_pdf(page_doc)
                page_doc.close()
                os.remove(page_pdf)
        pdf_merger.save(temp_searchable)
        pdf_merger.close()
        
        # Convert PDF to DOCX using pdf2docx monkey-patched methods
        convert_pdf_to_docx(temp_searchable, output_path)
        if os.path.exists(temp_searchable):
            os.remove(temp_searchable)
            
    elif output_type == 'text':
        full_text = []
        for idx, img_path in temp_files:
            txt_base = f"temp_ocr_out_p{idx}_{os.getpid()}"
            cmd = [tesseract_cmd] + extra_args + [img_path, txt_base]
            subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
            txt_file = txt_base + ".txt"
            if os.path.exists(txt_file):
                with open(txt_file, 'r', encoding='utf-8') as f:
                    full_text.append(clean_ocr(f.read()))
                os.remove(txt_file)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("\n\n--- PAGE BREAK ---\n\n".join(full_text))
            
    elif output_type == 'html':
        html_pages = []
        for idx, img_path in temp_files:
            hocr_base = f"temp_ocr_out_p{idx}_{os.getpid()}"
            cmd = [tesseract_cmd] + extra_args + [img_path, hocr_base, "hocr"]
            subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
            hocr_file = hocr_base + ".hocr"
            if os.path.exists(hocr_file):
                with open(hocr_file, 'r', encoding='utf-8') as f:
                    html_pages.append(f.read())
                os.remove(hocr_file)
        
        combined_html = "<html><head><title>OCR Reconstructed Document</title></head><body>"
        for page_content in html_pages:
            combined_html += f"<div class='ocr_page'>{page_content}</div><hr/>"
        combined_html += "</body></html>"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(combined_html)
            
    elif output_type == 'json':
        import json
        from PIL import Image
        json_data = {
            "confidence_average": 0.0,
            "total_words": 0,
            "pages": []
        }
        total_conf = 0.0
        word_count = 0
        
        for idx, img_path in temp_files:
            tsv_base = f"temp_ocr_out_p{idx}_{os.getpid()}"
            cmd = [tesseract_cmd] + extra_args + [img_path, tsv_base, "tsv"]
            subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
            tsv_file = tsv_base + ".tsv"
            
            page_words = []
            if os.path.exists(tsv_file):
                with open(tsv_file, 'r', encoding='utf-8') as f:
                    lines = f.readlines()
                
                headers = lines[0].strip().split('\t')
                for line in lines[1:]:
                    parts = line.strip().split('\t')
                    if len(parts) < len(headers):
                        continue
                    row = dict(zip(headers, parts))
                    text = row.get("text", "").strip()
                    conf = float(row.get("conf", "-1"))
                    
                    if text and conf >= 0:
                        page_words.append({
                            "text": text,
                            "confidence": conf,
                            "bbox": [
                                int(row.get("left", 0)),
                                int(row.get("top", 0)),
                                int(row.get("width", 0)),
                                int(row.get("height", 0))
                            ]
                        })
                        total_conf += conf
                        word_count += 1
                os.remove(tsv_file)
            
            page_width, page_height = 612, 792
            try:
                im = Image.open(img_path)
                page_width, page_height = im.size
            except Exception:
                pass
                
            json_data["pages"].append({
                "page_index": idx,
                "width": page_width,
                "height": page_height,
                "words": page_words
            })
        
        if word_count > 0:
            json_data["confidence_average"] = round(total_conf / word_count, 2)
        json_data["total_words"] = word_count
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(json_data, f, indent=2)

    # Clean up temp image files (only if we converted a PDF)
    if is_pdf:
        for _, img_path in temp_files:
            try:
                os.remove(img_path)
            except Exception:
                pass
    print("OCR task finished successfully.")

def compress_pdf(input_path, output_path, level, custom_dpi=None, custom_quality=None):
    print(f"Compressing PDF: {input_path} -> {output_path} (preset: {level})")
    doc = fitz.open(input_path)
    
    dpi = 200
    quality = 85
    if level == 'extreme':
        dpi = 96
        quality = 50
    elif level == 'high':
        dpi = 150
        quality = 70
    elif level == 'medium':  # 'recommended' maps to 'medium'
        dpi = 220
        quality = 82
    elif level == 'low':
        dpi = None
        quality = 95
    elif level == 'custom':
        dpi = int(custom_dpi) if custom_dpi else 200
        quality = int(custom_quality) if custom_quality else 85

    # Clear metadata to reduce overhead
    doc.set_metadata({})
    
    processed_images = {}
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        image_list = page.get_images(full=True)
        
        for img in image_list:
            xref = img[0]
            if xref in processed_images:
                continue
            
            try:
                base_image = doc.extract_image(xref)
                if not base_image:
                    continue
                image_bytes = base_image["image"]
                image_ext = base_image.get("ext", "").lower()
                is_lossless = image_ext in ("png", "gif", "tiff")
                
                from PIL import Image
                import io
                
                pil_img = Image.open(io.BytesIO(image_bytes))
                
                # Check resolution downsampling (only if dpi is set)
                if dpi:
                    rects = page.get_image_rects(xref)
                    if rects:
                        first_rect = rects[0]
                        img_rect = first_rect[0] if isinstance(first_rect, tuple) else first_rect
                        max_w = int(img_rect.width * (dpi / 72.0))
                        max_h = int(img_rect.height * (dpi / 72.0))
                    else:
                        page_rect = page.rect
                        max_w = int(page_rect.width * (dpi / 72.0))
                        max_h = int(page_rect.height * (dpi / 72.0))
                    
                    max_w = max(max_w, 32)
                    max_h = max(max_h, 32)
                    
                    w, h = pil_img.size
                    scale_factor = min(max_w / w, max_h / h)
                    
                    # Only downsample if scaling factor reduces size
                    if scale_factor < 1.0:
                        new_size = (int(w * scale_factor), int(h * scale_factor))
                        pil_img = pil_img.resize(new_size, Image.Resampling.LANCZOS)
                
                out_io = io.BytesIO()
                # Retain lossless format if originally lossless or containing transparency
                if is_lossless or pil_img.mode in ("RGBA", "LA") or (pil_img.mode == "P" and "transparency" in pil_img.info):
                    pil_img.save(out_io, format="PNG", optimize=True)
                else:
                    pil_img.save(out_io, format="JPEG", quality=quality, optimize=True)
                
                new_data = out_io.getvalue()
                # ONLY substitute if the compressed version is actually smaller
                if len(new_data) < len(image_bytes):
                    page.replace_image(xref, stream=new_data)
                processed_images[xref] = True
            except Exception as img_err:
                print(f"Skipped image compression for xref {xref}: {img_err}")
                
    doc.save(
        output_path,
        garbage=4,
        deflate=True,
        clean=True
    )
    doc.close()
    print("Compression complete.")

def edit_pdf(input_path, output_path, meta_path):
    print(f"Editing PDF: {input_path} -> {output_path} using meta {meta_path}")
    import json
    import base64
    
    with open(meta_path, "r", encoding="utf-8") as f:
        meta = json.load(f)
        
    doc = fitz.open(input_path)
    
    # 1. Page reorganization (pageOrder)
    page_order = meta.get("pageOrder")
    if page_order:
        new_doc = fitz.open()
        for idx in page_order:
            if idx == -1:
                # Add a blank page
                new_doc.new_page(width=612, height=792)
            elif 0 <= idx < len(doc):
                new_doc.insert_pdf(doc, from_page=idx, to_page=idx)
        doc.close()
        doc = new_doc
        
    # 2. Page rotations
    rotations = meta.get("rotations") or {}
    for page_idx_str, rot_angle in rotations.items():
        page_idx = int(page_idx_str)
        if 0 <= page_idx < len(doc):
            page = doc[page_idx]
            page.set_rotation((page.rotation + rot_angle) % 360)
            
    # Helper to convert hex colors to RGB floats (0.0 to 1.0)
    def hex_to_rgb(hex_str):
        if not hex_str:
            return (0, 0, 0)
        h = hex_str.lstrip("#")
        try:
            return tuple(int(h[i:i+2], 16) / 255.0 for i in (0, 2, 4))
        except Exception:
            return (0, 0, 0)
            
    # Helper to get font name
    def get_font_name(font_family, bold, italic):
        family = (font_family or "helvetica").lower()
        if "times" in family:
            if bold and italic: return "times-bolditalic"
            if bold: return "times-bold"
            if italic: return "times-italic"
            return "times-roman"
        elif "courier" in family:
            if bold and italic: return "courier-boldoblique"
            if bold: return "courier-bold"
            if italic: return "courier-oblique"
            return "courier"
        else:
            if bold and italic: return "helvetica-boldoblique"
            if bold: return "helvetica-bold"
            if italic: return "helvetica-oblique"
            return "helvetica"

    # 3. Process annotations
    annotations = meta.get("annotations") or []
    
    # Pre-pass: group text annotations with their mask shapes to run true text-replacement (redaction)
    locked_masks = {}
    for ann in annotations:
        if ann.get("type") == "shape" and ann.get("locked") and ann.get("bgColor") == "#ffffff":
            key = (ann.get("page"), round(ann.get("x"), 1), round(ann.get("y"), 1))
            locked_masks[key] = ann
            
    for ann in annotations:
        page_idx = ann.get("page")
        if not (0 <= page_idx < len(doc)):
            continue
            
        page = doc[page_idx]
        page_width = page.rect.width
        page_height = page.rect.height
        
        # Coordinates in PDF points (but from bottom-left coordinate system in frontend)
        # Convert to PyMuPDF's top-left coordinates:
        w_val = ann.get("width") or 50.0
        h_val = ann.get("height") or 30.0
        x_val = ann.get("x") or 0.0
        y_val = ann.get("y") or 0.0
        
        y0 = page_height - (y_val + h_val)
        y1 = page_height - y_val
        x0 = x_val
        x1 = x_val + w_val
        rect = fitz.Rect(x0, y0, x1, y1)
        
        ann_type = ann.get("type")
        
        # Check if this text annotation replaces existing text (true text editing)
        is_replacement = False
        if ann_type == "text" and ann.get("maskId"):
            is_replacement = True
        elif ann_type == "text":
            key = (page_idx, round(x_val, 1), round(y_val, 1))
            if key in locked_masks:
                is_replacement = True
                
        if is_replacement:
            # TRUE TEXT EDITING: Redact the original text region!
            page.add_redact_annot(rect, fill=(1, 1, 1))
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
            
        # Skip rendering locked masks since the redaction fill handles the white background
        if ann_type == "shape" and ann.get("locked"):
            continue
            
        if ann_type == "text":
            bg_color = ann.get("bgColor")
            if bg_color and bg_color != "transparent" and bg_color != "":
                page.draw_rect(rect, color=hex_to_rgb(bg_color), fill=hex_to_rgb(bg_color), width=0)
                
            text = ann.get("text") or ""
            font_size = ann.get("fontSize") or 12
            font_name = get_font_name(ann.get("fontFamily"), ann.get("bold"), ann.get("italic"))
            text_color = hex_to_rgb(ann.get("color"))
            
            align_map = {"left": 0, "center": 1, "right": 2}
            alignment = align_map.get(ann.get("alignment") or "left", 0)
            
            page.insert_textbox(
                rect,
                text,
                fontsize=font_size,
                fontname=font_name,
                color=text_color,
                align=alignment
            )
            
            # Underline / Strikethrough
            if ann.get("underline") or ann.get("strikethrough"):
                line_y = y0 + font_size * 0.95
                if ann.get("underline"):
                    page.draw_line(fitz.Point(x0, line_y), fitz.Point(x1, line_y), color=text_color, width=1)
                if ann.get("strikethrough"):
                    st_y = y0 + font_size * 0.5
                    page.draw_line(fitz.Point(x0, st_y), fitz.Point(x1, st_y), color=text_color, width=1)
                    
        elif ann_type == "shape":
            shape_type = ann.get("shapeType") or "rectangle"
            bg_color = ann.get("bgColor")
            fill_color = hex_to_rgb(bg_color) if (bg_color and bg_color != "transparent") else None
            border_color = hex_to_rgb(ann.get("color")) if ann.get("color") else (0, 0, 0)
            border_width = ann.get("borderWidth") if ann.get("borderWidth") is not None else 2
            
            if shape_type == "circle":
                page.draw_ellipse(rect, color=border_color, fill=fill_color, width=border_width)
            elif shape_type == "line":
                page.draw_line(fitz.Point(x0, y0), fitz.Point(x1, y1), color=border_color, width=border_width)
            else: # rectangle
                page.draw_rect(rect, color=border_color, fill=fill_color, width=border_width)
                
        elif ann_type == "drawing":
            color = hex_to_rgb(ann.get("color"))
            thickness = ann.get("borderWidth") if ann.get("borderWidth") is not None else 2
            paths = ann.get("paths")
            if paths:
                for path in paths:
                    pts = [fitz.Point(pt["x"], page_height - pt["y"]) for pt in path]
                    if len(pts) > 1:
                        page.draw_polyline(pts, color=color, width=thickness)
                        
        elif ann_type == "image":
            image_bytes = ann.get("imageBytes")
            if image_bytes:
                try:
                    if "," in image_bytes:
                        _, encoded = image_bytes.split(",", 1)
                    else:
                        encoded = image_bytes
                    img_data = base64.b64decode(encoded)
                    page.insert_image(rect, stream=img_data)
                except Exception as img_err:
                    print(f"Failed to embed image annotation: {img_err}")
                    
        elif ann_type == "highlight":
            bg_color = ann.get("bgColor")
            color = hex_to_rgb(bg_color) if bg_color else (1, 0.9, 0.2)
            annot = page.add_highlight_annot(rect)
            annot.set_colors(stroke=color)
            annot.update()
            
        elif ann_type == "underline":
            color = hex_to_rgb(ann.get("color")) or (0, 0, 1)
            annot = page.add_underline_annot(rect)
            annot.set_colors(stroke=color)
            annot.update()
            
        elif ann_type == "strikethrough":
            color = hex_to_rgb(ann.get("color")) or (1, 0, 0)
            annot = page.add_strikeout_annot(rect)
            annot.set_colors(stroke=color)
            annot.update()
            
        elif ann_type == "note":
            color = hex_to_rgb(ann.get("color")) or (0.93, 0.42, 0.3)
            note_content = ann.get("noteContent") or ""
            annot = page.add_text_annot(fitz.Point(x0, y0), note_content)
            annot.set_colors(stroke=color)
            annot.update()
            
        elif ann_type == "callout":
            bg_color = ann.get("bgColor")
            fill_color = hex_to_rgb(bg_color) if (bg_color and bg_color != "transparent") else (0.95, 0.95, 0.95)
            border_color = hex_to_rgb(ann.get("color")) or (0.93, 0.42, 0.3)
            border_width = ann.get("borderWidth") if ann.get("borderWidth") is not None else 1.5
            
            page.draw_rect(rect, color=border_color, fill=fill_color, width=border_width)
            
            text = ann.get("text") or ""
            font_size = ann.get("fontSize") or 10
            font_name = get_font_name(ann.get("fontFamily"), ann.get("bold"), ann.get("italic"))
            text_color = hex_to_rgb(ann.get("color"))
            
            page.insert_textbox(
                rect,
                text,
                fontsize=font_size,
                fontname=font_name,
                color=text_color,
                align=0
            )
            
    doc.save(output_path, garbage=3, deflate=True)
    doc.close()
    print("PDF editing and text replacement complete.")

def repair_pdf(input_path, output_path):
    print(f"Repairing PDF: {input_path} -> {output_path}")
    import json
    import shutil
    
    report = {
        "errors_found": [],
        "errors_repaired": [],
        "remaining_warnings": []
    }
    
    temp_patched = None
    input_path_to_open = input_path
    
    # Stage 1: Pre-Flight Binary Validation and Patching
    try:
        with open(input_path, 'rb') as f:
            data = f.read()
            
        original_len = len(data)
        modified = False
        
        # 1. Look for %PDF- signature. Trim any leading garbage bytes before it.
        pdf_start = data.find(b'%PDF-')
        if pdf_start > 0:
            data = data[pdf_start:]
            report["errors_found"].append(f"Found leading garbage bytes ({pdf_start} bytes) before PDF header signature.")
            report["errors_repaired"].append("Trimmed leading garbage bytes before PDF header signature.")
            modified = True
        elif pdf_start < 0:
            # Missing header completely! Prepend standard PDF-1.4 header descriptor
            data = b'%PDF-1.4\r\n' + data
            report["errors_found"].append("Invalid PDF Header signature (completely missing).")
            report["errors_repaired"].append("Injected standard '%PDF-1.4' header descriptor.")
            modified = True
            
        # 2. Look for %%EOF signature. Trim any excessive trailing junk or append it if missing.
        eof_idx = data.rfind(b'%%EOF')
        if eof_idx >= 0:
            trailing_len = len(data) - (eof_idx + 5)
            if trailing_len > 10:
                data = data[:eof_idx + 5]
                report["errors_found"].append(f"Trailing garbage bytes ({trailing_len} bytes) found after %%EOF marker.")
                report["errors_repaired"].append("Trimmed trailing garbage bytes after %%EOF marker.")
                modified = True
        else:
            # Completely missing EOF! Append it
            data = data + b'\r\n%%EOF\r\n'
            report["errors_found"].append("Missing EOF marker.")
            report["errors_repaired"].append("Appended standard EOF token block.")
            modified = True
            
        if modified:
            temp_patched = input_path + f".patched_{os.getpid()}"
            with open(temp_patched, 'wb') as f:
                f.write(data)
            input_path_to_open = temp_patched
            
    except Exception as val_err:
        report["errors_found"].append(f"Pre-flight binary analysis read error: {val_err}")
    
    # Stage 2 & 3: Structural Rebuild & Object Re-indexing
    fitz.TOOLS.mupdf_warnings()  # Clear warning buffer
    try:
        doc = fitz.open(input_path_to_open)
        load_warnings = fitz.TOOLS.mupdf_warnings()
        if load_warnings:
            for warn in load_warnings.split('\n'):
                warn = warn.strip()
                if warn:
                    report["errors_found"].append(warn)
                    report["errors_repaired"].append(f"Reconstructed malformed object: {warn}")
                    
        # Stage 4: Structure Optimization & Save
        doc.save(
            output_path,
            garbage=4,
            clean=True,
            deflate=True
        )
        doc.close()
        
        # Test loading repaired output
        test_doc = fitz.open(output_path)
        page_count = len(test_doc)
        test_doc.close()
        
        if page_count == 0:
            raise ValueError("Repaired document contains 0 pages.")
    except Exception as rep_err:
        report["errors_found"].append(f"Fatal structural recovery exception: {rep_err}")
        # Fall back to copy
        try:
            shutil.copyfile(input_path, output_path)
        except Exception:
            pass
        report["remaining_warnings"].append(f"Multi-stage recovery failed to reconstruct structures: {rep_err}")
    finally:
        if temp_patched and os.path.exists(temp_patched):
            try:
                os.remove(temp_patched)
            except Exception:
                pass
        
    if not report["errors_found"]:
        report["errors_repaired"].append("Verified cross-reference table and rebuilt structures successfully.")
        
    # Write report file alongside repaired output
    report_path = output_path + ".report.json"
    with open(report_path, 'w', encoding='utf-8') as rf:
        json.dump(report, rf, indent=2)
    print("Repair complete.")

def detect_form_fields(input_path, output_path):
    print(f"Detecting form fields: {input_path} -> {output_path}")
    import json
    doc = fitz.open(input_path)
    fields = []
    
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_width = page.rect.width
        page_height = page.rect.height
        
        for widget in page.widgets():
            ftype = widget.field_type
            rect = widget.rect
            
            # Map layout rect to page percentages
            x = (rect.x0 / page_width) * 100
            y = (rect.y0 / page_height) * 100
            w = ((rect.x1 - rect.x0) / page_width) * 100
            h = ((rect.y1 - rect.y0) / page_height) * 100
            
            field_type = "text"
            options = []
            
            if ftype == 1:
                field_type = "signature"
            elif ftype == 2:
                field_type = "checkbox"
                # Some checkbox fields represent radio buttons in Acrobat PDF spec
                if widget.is_commit() or "radio" in (widget.field_name or "").lower():
                    field_type = "radio"
            elif ftype == 3:
                field_type = "dropdown"
                options = widget.choice_values() if hasattr(widget, "choice_values") else []
            else:
                field_type = "text"
                if "date" in (widget.field_name or "").lower():
                    field_type = "date"
                    
            fields.append({
                "id": widget.field_name or f"field_{page_idx}_{len(fields)}",
                "page": page_idx,
                "type": field_type,
                "x": x,
                "y": y,
                "width": w,
                "height": h,
                "name": widget.field_name or f"field_{page_idx}_{len(fields)}",
                "value": widget.field_value or "",
                "options": options,
                "required": widget.is_required() if hasattr(widget, "is_required") else False
            })
            
    doc.close()
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(fields, f, indent=2)
    print("Form field detection complete.")

def compare_pdfs(pdf1_path, pdf2_path, output_diff_pdf, output_report_json, output_report_html):
    print(f"Comparing PDFs: {pdf1_path} vs {pdf2_path}")
    import json
    import difflib
    from PIL import Image, ImageChops
    
    doc1 = fitz.open(pdf1_path)
    doc2 = fitz.open(pdf2_path)
    
    report = {
        "summary": {
            "total_differences": 0,
            "added": 0,
            "removed": 0,
            "modified": 0
        },
        "differences": []
    }
    
    # Open visual copies to highlight modifications
    highlight_doc1 = fitz.open(pdf1_path)
    highlight_doc2 = fitz.open(pdf2_path)
    
    max_pages = max(len(doc1), len(doc2))
    for page_idx in range(max_pages):
        if page_idx < len(doc1) and page_idx < len(doc2):
            p1 = doc1[page_idx]
            p2 = doc2[page_idx]
            
            words1 = p1.get_text("words")
            words2 = p2.get_text("words")
            
            text1_list = [w[4] for w in words1]
            text2_list = [w[4] for w in words2]
            
            sm = difflib.SequenceMatcher(None, text1_list, text2_list)
            opcodes = sm.get_opcodes()
            
            hp1 = highlight_doc1[page_idx]
            hp2 = highlight_doc2[page_idx]
            
            for tag, i1, i2, j1, j2 in opcodes:
                if tag == 'replace':
                    report["summary"]["modified"] += 1
                    for idx in range(i1, i2):
                        w = words1[idx]
                        rect = fitz.Rect(w[0], w[1], w[2], w[3])
                        hp1.add_highlight_annot(rect).set_colors(stroke=(1.0, 0.9, 0.2))
                        report["differences"].append({
                            "page": page_idx,
                            "type": "text_modified_removed",
                            "content": w[4],
                            "bbox": [w[0], w[1], w[2], w[3]]
                        })
                    for idx in range(j1, j2):
                        w = words2[idx]
                        rect = fitz.Rect(w[0], w[1], w[2], w[3])
                        hp2.add_highlight_annot(rect).set_colors(stroke=(1.0, 0.9, 0.2))
                        report["differences"].append({
                            "page": page_idx,
                            "type": "text_modified_added",
                            "content": w[4],
                            "bbox": [w[0], w[1], w[2], w[3]]
                        })
                elif tag == 'delete':
                    report["summary"]["removed"] += 1
                    for idx in range(i1, i2):
                        w = words1[idx]
                        rect = fitz.Rect(w[0], w[1], w[2], w[3])
                        hp1.add_highlight_annot(rect).set_colors(stroke=(1.0, 0.2, 0.2))
                        # Highlight deletion coordinates on modified page as a red rectangle with a light red fill
                        hp2.add_rect_annot(rect).set_colors(stroke=(1.0, 0.2, 0.2), fill=(1.0, 0.85, 0.85))
                        report["differences"].append({
                            "page": page_idx,
                            "type": "text_removed",
                            "content": w[4],
                            "bbox": [w[0], w[1], w[2], w[3]]
                        })
                elif tag == 'insert':
                    report["summary"]["added"] += 1
                    for idx in range(j1, j2):
                        w = words2[idx]
                        rect = fitz.Rect(w[0], w[1], w[2], w[3])
                        hp2.add_highlight_annot(rect).set_colors(stroke=(0.2, 1.0, 0.2))
                        report["differences"].append({
                            "page": page_idx,
                            "type": "text_added",
                            "content": w[4],
                            "bbox": [w[0], w[1], w[2], w[3]]
                        })
            
            # Visual layout pixel comparison
            try:
                pix1 = p1.get_pixmap(dpi=72)
                pix2 = p2.get_pixmap(dpi=72)
                if pix1.width == pix2.width and pix1.height == pix2.height:
                    img1 = Image.frombytes("RGB", [pix1.width, pix1.height], pix1.samples)
                    img2 = Image.frombytes("RGB", [pix2.width, pix2.height], pix2.samples)
                    diff = ImageChops.difference(img1, img2)
                    diff_box = diff.getbbox()
                    if diff_box:
                        rect = fitz.Rect(diff_box[0], diff_box[1], diff_box[2], diff_box[3])
                        # Filter out overlaps with text diffs to prevent clutter
                        hp2.add_rect_annot(rect).set_colors(stroke=(1.0, 0.5, 0.0))
                        report["differences"].append({
                            "page": page_idx,
                            "type": "visual_difference",
                            "content": "Visual shift / logo or image change detected.",
                            "bbox": [rect.x0, rect.y0, rect.x1, rect.y1]
                        })
                        report["summary"]["total_differences"] += 1
            except Exception as vis_err:
                print(f"Visual compare skipped: {vis_err}")
                
        elif page_idx < len(doc1):
            report["summary"]["removed"] += 1
            report["differences"].append({
                "page": page_idx,
                "type": "page_removed",
                "content": f"Page {page_idx + 1} completely removed."
            })
        else:
            report["summary"]["added"] += 1
            report["differences"].append({
                "page": page_idx,
                "type": "page_added",
                "content": f"Page {page_idx + 1} completely added."
            })
            
    report["summary"]["total_differences"] += (
        report["summary"]["added"] + report["summary"]["removed"] + report["summary"]["modified"]
    )
    
    # Save highlighted modified PDF as the main output
    highlight_doc2.save(output_diff_pdf, garbage=4, deflate=True)
    
    # Save highlighted original PDF as the extra file (.original.pdf)
    highlight_doc1_path = output_diff_pdf + ".original.pdf"
    highlight_doc1.save(highlight_doc1_path, garbage=4, deflate=True)
    
    highlight_doc1.close()
    highlight_doc2.close()
    doc1.close()
    doc2.close()
    
    # Save reports
    with open(output_report_json, 'w', encoding='utf-8') as jf:
        json.dump(report, jf, indent=2)
        
    html_content = f"""<html>
<head>
    <title>PDF Comparison Report</title>
    <style>
        body {{ font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; margin: 40px; background-color: #f8fafc; color: #1e293b; }}
        h1 {{ color: #0f172a; border-bottom: 2px solid #e2e8f0; padding-bottom: 10px; }}
        .summary-card {{ display: flex; gap: 20px; margin-bottom: 30px; }}
        .stat-box {{ flex: 1; padding: 20px; border-radius: 12px; background: white; box-shadow: 0 4px 6px -1px rgba(0,0,0,0.05); text-align: center; border: 1px solid #e2e8f0; }}
        .stat-num {{ font-size: 2rem; font-weight: bold; margin-bottom: 5px; }}
        .stat-label {{ font-size: 0.875rem; color: #64748b; text-transform: uppercase; letter-spacing: 0.5px; }}
        .diff-list {{ background: white; border-radius: 12px; padding: 20px; box-shadow: 0 4px 6px -1px rgba(0,0,0,0.05); border: 1px solid #e2e8f0; }}
        .diff-item {{ padding: 12px; border-bottom: 1px solid #f1f5f9; display: flex; justify-content: space-between; align-items: center; }}
        .diff-item:last-child {{ border-bottom: none; }}
        .badge {{ padding: 4px 10px; border-radius: 20px; font-size: 0.75rem; font-weight: 600; text-transform: uppercase; }}
        .badge-added {{ background-color: #dcfce7; color: #15803d; }}
        .badge-removed {{ background-color: #fee2e2; color: #b91c1c; }}
        .badge-modified {{ background-color: #fef9c3; color: #a16207; }}
        .badge-visual {{ background-color: #f0fdf4; color: #166534; }}
    </style>
</head>
<body>
    <h1>Document Comparison Report</h1>
    <div class="summary-card">
        <div class="stat-box" style="border-top: 4px solid #3b82f6;">
            <div class="stat-num">{report["summary"]["total_differences"]}</div>
            <div class="stat-label">Total Differences</div>
        </div>
        <div class="stat-box" style="border-top: 4px solid #22c55e;">
            <div class="stat-num" style="color: #22c55e;">{report["summary"]["added"]}</div>
            <div class="stat-label">Added Items</div>
        </div>
        <div class="stat-box" style="border-top: 4px solid #ef4444;">
            <div class="stat-num" style="color: #ef4444;">{report["summary"]["removed"]}</div>
            <div class="stat-label">Removed Items</div>
        </div>
        <div class="stat-box" style="border-top: 4px solid #eab308;">
            <div class="stat-num" style="color: #eab308;">{report["summary"]["modified"]}</div>
            <div class="stat-label">Modified Items</div>
        </div>
    </div>
    
    <h2>Change Details Log</h2>
    <div class="diff-list">
"""
    for diff in report["differences"]:
        badge_class = "badge-modified"
        if "added" in diff["type"]:
            badge_class = "badge-added"
        elif "removed" in diff["type"]:
            badge_class = "badge-removed"
        elif "visual" in diff["type"]:
            badge_class = "badge-visual"
            
        html_content += f"""
        <div class="diff-item">
            <div>
                <strong>Page {diff["page"] + 1}:</strong> 
                <span style="margin-left: 10px; font-family: monospace;">"{diff.get("content", "")}"</span>
            </div>
            <span class="badge {badge_class}">{diff["type"].replace('text_', '')}</span>
        </div>
"""
    html_content += """
    </div>
</body>
</html>
"""
    with open(output_report_html, 'w', encoding='utf-8') as hf:
        hf.write(html_content)
    print("Comparison complete.")

def preprocess_excel(input_path, output_path):
    print(f"Preprocessing Excel file: {input_path} -> {output_path}")
    import pandas as pd
    import openpyxl
    from openpyxl.styles import Alignment
    from openpyxl.utils import get_column_letter
    import re

    # 1. Load the xls/xlsx file
    if input_path.lower().endswith('.xls'):
        # Load all sheets
        sheets = pd.read_excel(input_path, sheet_name=None)
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            for sheet_name, df in sheets.items():
                df.to_excel(writer, sheet_name=sheet_name, index=False)
    else:
        # Copy / load as is
        import shutil
        shutil.copyfile(input_path, output_path)

    # 2. Modify properties in openpyxl
    wb = openpyxl.load_workbook(output_path)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        
        # Enable Fit to Page Width
        ws.sheet_properties.pageSetUpPr.fitToPage = True
        ws.page_setup.fitToWidth = 1
        ws.page_setup.fitToHeight = 0 # let it flow vertically
        
        # Set Orientation to Landscape for wider tables
        ws.page_setup.orientation = ws.ORIENTATION_LANDSCAPE
        
        # Configure columns and cell word wrap
        for col in ws.columns:
            max_len = 0
            for cell in col:
                if cell.value is not None:
                    max_len = max(max_len, len(str(cell.value)))
                    # Format numbers right with an indent of 1 to prevent text clipping/running into next cell
                    is_numeric = isinstance(cell.value, (int, float))
                    if not is_numeric and isinstance(cell.value, str):
                        cleaned = cell.value.strip().replace(',', '')
                        if re.match(r'^-?\d+(\.\d+)?$', cleaned):
                            is_numeric = True
                            
                    if is_numeric:
                        cell.alignment = Alignment(wrap_text=True, vertical='center', horizontal='right', indent=1)
                    else:
                        cell.alignment = Alignment(wrap_text=True, vertical='center', horizontal='left')
            
            col_letter = get_column_letter(col[0].column)
            # Increase column width spacing to avoid clipping
            ws.column_dimensions[col_letter].width = min(max(max_len + 5, 12), 30)

    wb.save(output_path)
    print("Preprocessed Excel file successfully.")

def convert_office_to_pdf(input_path, output_path):
    import os
    import sys
    
    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)
    
    ext = os.path.splitext(input_path.lower())[1]
    print(f"Converting Office document to PDF using native COM automation: {input_path} -> {output_path}")
    
    if ext in ['.docx', '.doc']:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        word = win32com.client.DispatchEx('Word.Application')
        word.Visible = False
        word.DisplayAlerts = False
        doc = None
        try:
            doc = word.Documents.Open(input_path)
            doc.SaveAs(output_path, FileFormat=17) # wdFormatPDF
            doc.Close()
            print("Successfully converted Word document to PDF.")
        finally:
            word.Quit()
            
    elif ext in ['.pptx', '.ppt']:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        powerpoint = win32com.client.DispatchEx('PowerPoint.Application')
        pres = None
        try:
            pres = powerpoint.Presentations.Open(input_path, WithWindow=False)
            pres.SaveAs(output_path, FileFormat=32) # ppSaveAsPDF
            pres.Close()
            print("Successfully converted PowerPoint presentation to PDF.")
        finally:
            powerpoint.Quit()
            
    elif ext in ['.xlsx', '.xls']:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        excel = win32com.client.DispatchEx('Excel.Application')
        excel.Visible = False
        excel.DisplayAlerts = False
        wb = None
        try:
            wb = excel.Workbooks.Open(input_path)
            wb.ExportAsFixedFormat(0, output_path) # xlTypePDF
            wb.Close(SaveChanges=False)
            print("Successfully converted Excel sheet to PDF.")
        finally:
            excel.Quit()
            
    else:
        raise ValueError(f"Unsupported office file extension: {ext}")

def main():
    parser = argparse.ArgumentParser(description="DocRIt Python Conversion Worker")
    parser.add_argument("--task", required=True, choices=["pdf-to-docx", "pdf-to-xlsx", "pdf-to-pptx", "ocr", "preprocess-excel", "office-to-pdf", "compress", "repair", "detect-forms", "compare", "edit"])
    parser.add_argument("--input", required=True, help="Path to input document")
    parser.add_argument("--output", required=True, help="Path to write converted output document")
    parser.add_argument("--ocr-type", default="text", choices=["text", "pdf", "docx", "html", "json"], help="OCR output format")
    
    # Extra arguments for advanced tasks
    parser.add_argument("--level", default="recommended", choices=["low", "medium", "high", "extreme", "custom"], help="Compression level preset")
    parser.add_argument("--dpi", type=int, help="Custom downsampling DPI")
    parser.add_argument("--quality", type=int, help="Custom JPEG quality")
    parser.add_argument("--input2", help="Secondary PDF path for document comparison")
    parser.add_argument("--report-json", help="Path to write compare JSON report")
    parser.add_argument("--report-html", help="Path to write compare HTML report")
    parser.add_argument("--extra", help="Path to JSON metadata payload (e.g. annotations for edit task)")
    
    args = parser.parse_args()
    
    if not os.path.exists(args.input):
        print(f"Error: Input file {args.input} does not exist.", file=sys.stderr)
        sys.exit(1)
        
    try:
        if args.task == "pdf-to-docx":
            convert_pdf_to_docx(args.input, args.output)
        elif args.task == "pdf-to-xlsx":
            convert_pdf_to_xlsx(args.input, args.output)
        elif args.task == "pdf-to-pptx":
            convert_pdf_to_pptx(args.input, args.output)
        elif args.task == "ocr":
            perform_ocr(args.input, args.output, args.ocr_type)
        elif args.task == "preprocess-excel":
            preprocess_excel(args.input, args.output)
        elif args.task == "office-to-pdf":
            convert_office_to_pdf(args.input, args.output)
        elif args.task == "compress":
            compress_pdf(args.input, args.output, args.level, args.dpi, args.quality)
        elif args.task == "repair":
            repair_pdf(args.input, args.output)
        elif args.task == "detect-forms":
            detect_form_fields(args.input, args.output)
        elif args.task == "edit":
            if not args.extra:
                raise ValueError("Edit task requires --extra argument with JSON path.")
            edit_pdf(args.input, args.output, args.extra)
        elif args.task == "compare":
            if not args.input2 or not args.report_json or not args.report_html:
                raise ValueError("Compare task requires --input2, --report-json, and --report-html arguments.")
            compare_pdfs(args.input, args.input2, args.output, args.report_json, args.report_html)
        sys.exit(0)
    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"Worker Error: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
